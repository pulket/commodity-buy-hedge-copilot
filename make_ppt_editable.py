"""
Builds a FULLY EDITABLE PowerPoint of the Commodity Copilot slides.
Every element is a native PowerPoint shape / text box / table (no images) so it
can be edited in PowerPoint. Numbers are pulled live from copilot_results.json.
"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR

HERE = os.path.dirname(os.path.abspath(__file__))
d = json.load(open(os.path.join(HERE, "copilot_results.json")))
C = d["commodities"]; port = d["portfolio"]; cfg = d["config"]
flag = next(n for n, r in C.items() if r["commodity"]["flagship"])
P = C[flag]; opt = P["optimise"]; ch = opt["chosen"]; sp = opt["split"]
sup = P["supply_chain"]; dist = opt["dist"]; act = P["act"]; bt = P["backtest"]

# ---- colours ----
PURPLE = RGBColor(0x78, 0x23, 0xDC); PURPLE_D = RGBColor(0x4C, 0x1D, 0x95)
TEAL = RGBColor(0x0F, 0x76, 0x6E); INK = RGBColor(0x20, 0x22, 0x2B)
MUTE = RGBColor(0x6B, 0x72, 0x80); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT = RGBColor(0xF5, 0xEE, 0xFD); TINT2 = RGBColor(0xEB, 0xDF, 0xFB)
TEALT = RGBColor(0xE6, 0xF6, 0xF3); LINE = RGBColor(0xD8, 0xD5, 0xE4)
GREEN = RGBColor(0x15, 0x80, 0x3D); RED = RGBColor(0xC0, 0x24, 0x3B)
AMBER = RGBColor(0xB4, 0x53, 0x09); PAPER = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
FONT = "Arial"


# ---- helpers ----
def rect(sl, x, y, w, h, fill=WHITE, line=LINE, lw=0.75, rounded=True, shadow=False):
    shp = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try: shp.adjustments[0] = 0.08
        except Exception: pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def txt(sl, x, y, w, h, runs, size=10, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, {})]
    p = tf.paragraphs[0]; p.alignment = align
    try: p.line_spacing = spacing
    except Exception: pass
    for t, st in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(st.get("size", size)); r.font.name = FONT
        r.font.bold = st.get("bold", bold)
        r.font.color.rgb = st.get("color", color)
        if st.get("italic"): r.font.italic = True
    return tb


def multiline(sl, x, y, w, h, lines, size=9, color=INK, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(3); tf.margin_right = Pt(3); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    tf.vertical_anchor = anchor
    for i, (t, st) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = st.get("align", align)
        try: p.line_spacing = st.get("sp", 1.0)
        except Exception: pass
        r = p.add_run(); r.text = t
        r.font.size = Pt(st.get("size", size)); r.font.name = FONT
        r.font.bold = st.get("bold", False)
        r.font.color.rgb = st.get("color", color)
        if st.get("italic"): r.font.italic = True
    return tb


def chip_row(sl, x, y, w, h, boxes, gap=0.06):
    n = len(boxes); bw = (w - gap * (n - 1)) / n
    for i, (title, sub, fill, tcol) in enumerate(boxes):
        bx = x + i * (bw + gap)
        rect(sl, bx, y, bw, h, fill=fill, line=(PURPLE if fill in (TINT, TINT2) else (TEAL if fill == TEALT else LINE)))
        multiline(sl, bx, y, bw, h, [
            (title, {"size": 7.6, "bold": True, "color": tcol, "sp": 0.95}),
            (sub, {"size": 6.4, "color": MUTE, "sp": 0.95})], anchor=MSO_ANCHOR.MIDDLE)


def footer(sl, page):
    rect(sl, 0.35, 7.02, 12.63, 0.005, fill=LINE, line=None, rounded=False)
    txt(sl, 0.35, 7.06, 6, 0.3, [("kearney", {"bold": True, "color": PURPLE, "size": 12}),
                                 ("  |  AI Commodity Copilot", {"color": INK, "size": 10, "bold": True})])
    txt(sl, 5.0, 7.12, 6.5, 0.3, "Illustrative model output - not financial advice.",
        size=7.5, color=MUTE, align=PP_ALIGN.CENTER)
    txt(sl, 12.5, 7.1, 0.45, 0.3, str(page), size=8, color=MUTE, align=PP_ALIGN.RIGHT)


def header(sl, eyebrow, title_runs):
    rect(sl, 0, 0, 13.333, 0.07, fill=PURPLE, line=None, rounded=False)
    txt(sl, 0.35, 0.13, 12, 0.25, eyebrow, size=9.5, color=PURPLE, bold=True)
    txt(sl, 0.35, 0.34, 12.7, 0.72, title_runs, size=17, color=INK, bold=True, spacing=1.05)
    rect(sl, 0.35, 1.08, 12.63, 0.02, fill=INK, line=None, rounded=False)


# ============================================================ SLIDE 1
s1 = prs.slides.add_slide(BLANK)
header(s1, "AI COMMODITY COPILOT  ·  SOLUTION & ARCHITECTURE",
       [("An ", {}), ("agentic Sense→Score→Decide→Act", {"color": PURPLE}),
        (" copilot simulates thousands of price futures, picks the buy/hedge mix best on ", {}),
        ("cost and stability", {"color": PURPLE}), (", and acts inside the council's guardrails", {})])

# ---- left column ----
LX, LW = 0.35, 3.45
# how it decides
rect(s1, LX, 1.2, LW, 1.15)
txt(s1, LX + 0.05, 1.24, LW, 0.2, "HOW IT DECIDES", size=9, bold=True, color=PURPLE_D)
multiline(s1, LX + 0.05, 1.44, LW - 0.1, 0.9, [
    ("1  Sense 6 signals per commodity (price, curve, FX, freight)", {"size": 8.3, "sp": 1.05}),
    ("2  Score a 0-100 buy-now signal → buy / phase / wait", {"size": 8.3, "sp": 1.05}),
    ("3  Decide via Monte-Carlo: minimise RALC on the frontier", {"size": 8.3, "sp": 1.05}),
    ("4  Act auto-execute in caps, else escalate to the council", {"size": 8.3, "sp": 1.05})],
    anchor=MSO_ANCHOR.TOP)
# metric
rect(s1, LX, 2.45, LW, 1.0, fill=TINT)
txt(s1, LX + 0.05, 2.49, LW, 0.2, "THE METRIC (COST + STABILITY)", size=9, bold=True, color=PURPLE_D)
multiline(s1, LX + 0.05, 2.69, LW - 0.1, 0.75, [
    ("RALC = E[Cost] + λ·SD[Cost]", {"size": 9.5, "bold": True, "color": PURPLE_D, "sp": 1.0}),
    ("Minimise expected cost plus a penalty on how much it can swing. λ set by 3 presets. "
     + f"{cfg['n_paths']:,} MC paths.", {"size": 7.6, "color": MUTE, "sp": 1.02})], anchor=MSO_ANCHOR.TOP)
# data
rect(s1, LX, 3.55, LW, 0.95)
txt(s1, LX + 0.05, 3.59, LW, 0.2, "DATA (THIS FOLDER)", size=9, bold=True, color=PURPLE_D)
multiline(s1, LX + 0.05, 3.78, LW - 0.1, 0.7, [
    ("Financial: palm/crude/silver 1/3/6-mo, FX (USDMYR, USDINR), Baltic Dry freight.",
     {"size": 7.5, "sp": 1.05}),
    ("Supply-chain: real demand & volatility from the sales file, inventory, suppliers.",
     {"size": 7.5, "color": TEAL, "sp": 1.05})], anchor=MSO_ANCHOR.TOP)
# presets table
rect(s1, LX, 4.6, LW, 1.6)
txt(s1, LX + 0.05, 4.64, LW, 0.2, "3 RISK-APPETITE PRESETS (COUNCIL PICKS)", size=8.5, bold=True, color=PURPLE_D)
ptab = s1.shapes.add_table(4, 5, Inches(LX + 0.05), Inches(4.86), Inches(LW - 0.1), Inches(1.2)).table
ptab.columns[0].width = Inches(1.15)
for _cj in range(1, 5):
    ptab.columns[_cj].width = Inches((LW - 0.1 - 1.15) / 4)
hdr = ["Preset", "λ", "Cover", "Hedge", "Auto"]
rowsp = [("Conservative", "1.5", "85%", "80%", "₹60Cr"),
         ("Balanced", "0.8", "60%", "50%", "₹40Cr"),
         ("Aggressive", "0.3", "40%", "30%", "₹20Cr")]
for j, hh in enumerate(hdr):
    cel = ptab.cell(0, j); cel.text = hh
    cel.fill.solid(); cel.fill.fore_color.rgb = TINT
for i, rr in enumerate(rowsp, 1):
    for j, v in enumerate(rr):
        cel = ptab.cell(i, j); cel.text = v
        cel.fill.solid(); cel.fill.fore_color.rgb = WHITE
for i in range(4):
    for j in range(5):
        c = ptab.cell(i, j); c.margin_left = Pt(3); c.margin_right = Pt(2)
        c.margin_top = Pt(1); c.margin_bottom = Pt(1); c.vertical_anchor = MSO_ANCHOR.MIDDLE
        for pr in c.text_frame.paragraphs:
            for r in pr.runs:
                r.font.size = Pt(7.6); r.font.name = FONT
                r.font.color.rgb = PURPLE_D if i == 0 else INK
                r.font.bold = (i == 0 or j == 0)

# ---- architecture ----
AX, AW = 4.0, 8.98
rect(s1, AX, 1.2, AW, 5.0)
txt(s1, AX + 0.1, 1.26, 6, 0.22, "ARCHITECTURE - TWO TRACKS RUN THROUGH EVERY STAGE",
    size=9, bold=True, color=PURPLE_D)
# legend
rect(s1, AX + 6.55, 1.24, 0.32, 0.16, fill=PURPLE_D, line=None); txt(s1, AX + 6.5, 1.24, 0.42, 0.16, "FIN", size=6.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s1, AX + 6.9, 1.24, 0.7, 0.16, "financial", size=7, color=MUTE, anchor=MSO_ANCHOR.MIDDLE)
rect(s1, AX + 7.55, 1.24, 0.3, 0.16, fill=TEAL, line=None); txt(s1, AX + 7.5, 1.24, 0.4, 0.16, "SC", size=6.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s1, AX + 7.88, 1.24, 1.0, 0.16, "supply-chain", size=7, color=MUTE, anchor=MSO_ANCHOR.MIDDLE)

# AI rail
RAILX = AX + 7.35
rect(s1, RAILX, 1.55, 1.5, 4.55, fill=TINT)
txt(s1, RAILX + 0.06, 1.6, 1.4, 0.3, "AI / LLM LAYER (WRAPS ALL)", size=7.6, bold=True, color=PURPLE_D)
ai = [("Sense - LLM", "reads news/MPOB/USDA → cited signal, extracts ERP data"),
      ("Decide - ML", "path probabilities + anomaly alerts"),
      ("Act - LLM agent", "drafts & executes in caps, writes escalation memo"),
      ("NL Q&A - LLM", "instant what-if + explanations"),
      ("Score = rules", "kept auditable; AI never sets the price")]
ay = 1.95
for t, sub in ai:
    multiline(s1, RAILX + 0.06, ay, 1.4, 0.78, [
        (t, {"size": 7.4, "bold": True, "color": PURPLE_D, "sp": 0.95}),
        (sub, {"size": 6.4, "color": MUTE, "sp": 0.95})], anchor=MSO_ANCHOR.TOP)
    ay += 0.82

# stages
STAGE_X = AX + 0.12; TAG_W = 0.8
LANE_X = STAGE_X + TAG_W + 0.08; LANE_W = RAILX - LANE_X - 0.12
stages = [
    ("1 · SENSE", PURPLE_D,
     [("Price & curve", "%ile, momentum", WHITE), ("Volatility", "risk", WHITE), ("FX", "USDMYR/INR", WHITE), ("Freight", "Baltic Dry", WHITE)],
     [("Demand", "sales file", TEALT), ("Inventory", "on hand", TEALT), ("Lead time", "order→arrive", TEALT), ("Suppliers", "capacity", TEALT)]),
    ("2 · SCORE", PURPLE_D,
     [("Value", "cheap?", TINT), ("Momentum", "rising?", TINT), ("Curve", "contango?", TINT), ("FX", "₹ weaker?", TINT), ("→ Buy-now", "buy/phase/wait", TINT)],
     [("Demand vol", f"CV {sup['demand_cv']*100:.0f}%", TEALT), ("Cover", f"{sup['months_cover']:.1f} mo", TEALT), ("Reorder", "breached" if sup['below_reorder'] else "ok", TEALT), ("→ Urgency", "continuity", TEALT)]),
    ("3 · DECIDE", PURPLE_D,
     [("Cover %", "lock now", TINT2), ("Hedge %", "options/caps", TINT2), ("RALC min", "on frontier", TINT2), ("Triggers", "price/vol/FX", TINT2)],
     [("Net buy", f"{sup['net_procurement']:,.0f}", TEALT), ("Must-cover", f"{sup['must_cover_now']:,.0f} now", TEALT), ("Supplier cap", f"{sup['supplier_share']*100:.0f}%", TEALT), ("Safety stock", "from CV", TEALT)]),
]
STAGE_Y0 = 1.62; STAGE_H = 0.88
for i, (tag, tagcol, fin, sc) in enumerate(stages):
    y0 = STAGE_Y0 + i * STAGE_H
    rect(s1, STAGE_X, y0, TAG_W, 0.74, fill=tagcol, line=None)
    multiline(s1, STAGE_X, y0, TAG_W, 0.74, [(tag, {"size": 7.6, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "sp": 0.95})], anchor=MSO_ANCHOR.MIDDLE)
    # FIN lane
    rect(s1, LANE_X - 0.02, y0, 0.2, 0.34, fill=PURPLE_D, line=None)
    txt(s1, LANE_X - 0.06, y0, 0.28, 0.34, "FIN", size=5.6, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    chip_row(s1, LANE_X + 0.22, y0, LANE_W - 0.22, 0.34, [(t, s, f, INK if f != WHITE else INK) for t, s, f in fin])
    # SC lane
    rect(s1, LANE_X - 0.02, y0 + 0.38, 0.2, 0.34, fill=TEAL, line=None)
    txt(s1, LANE_X - 0.06, y0 + 0.38, 0.28, 0.34, "SC", size=5.6, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    chip_row(s1, LANE_X + 0.22, y0 + 0.38, LANE_W - 0.22, 0.34, [(t, s, TEALT, TEAL) for t, s, f in sc])

# ACT + OUTPUT (wide lanes)
def wide_stage(i, tag, tagcol, fin_txt, sc_txt):
    y0 = STAGE_Y0 + i * STAGE_H
    rect(s1, STAGE_X, y0, TAG_W, 0.74, fill=tagcol, line=None)
    multiline(s1, STAGE_X, y0, TAG_W, 0.74, [(tag, {"size": 7.6, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "sp": 0.95})], anchor=MSO_ANCHOR.MIDDLE)
    rect(s1, LANE_X - 0.02, y0, 0.2, 0.34, fill=PURPLE_D, line=None)
    txt(s1, LANE_X - 0.06, y0, 0.28, 0.34, "FIN", size=5.6, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s1, LANE_X + 0.22, y0, LANE_W - 0.22, 0.34, fill=TINT, line=PURPLE)
    multiline(s1, LANE_X + 0.28, y0, LANE_W - 0.3, 0.34, fin_txt, anchor=MSO_ANCHOR.MIDDLE)
    rect(s1, LANE_X - 0.02, y0 + 0.38, 0.2, 0.34, fill=TEAL, line=None)
    txt(s1, LANE_X - 0.06, y0 + 0.38, 0.28, 0.34, "SC", size=5.6, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s1, LANE_X + 0.22, y0 + 0.38, LANE_W - 0.22, 0.34, fill=TEALT, line=TEAL)
    multiline(s1, LANE_X + 0.28, y0 + 0.38, LANE_W - 0.3, 0.34, sc_txt, anchor=MSO_ANCHOR.MIDDLE)

wide_stage(3, "4 · ACT", PURPLE,
           [("Auto-execute inside caps ", {"size": 7.8, "bold": True, "color": PURPLE_D}), ("- else escalate to council, full audit trail, manual override", {"size": 7, "color": MUTE})],
           [("Supply continuity first ", {"size": 7.8, "bold": True, "color": TEAL}), ("- mandatory must-cover buy before price-optimising", {"size": 7, "color": MUTE})])
wide_stage(4, "5 · OUTPUT", PURPLE_D,
           [(f"Hedge plan + portfolio nets risk → {port['diversification_benefit_pct']:.0f}% natural hedge ", {"size": 7.8, "bold": True, "color": PURPLE_D}), ("→ Treasury", {"size": 7, "color": MUTE})],
           [("Buy plan: volumes & timing ", {"size": 7.8, "bold": True, "color": TEAL}), ("→ Procurement", {"size": 7, "color": MUTE})])

# bottom band
band = [("Two tracks", "Financial + Supply-chain, every stage"),
        ("Cost + stability", "RALC on an efficient frontier"),
        ("Governed autonomy", "caps, escalation, audit, override"),
        ("Backtested", "steadier cost vs naive on real history")]
bw = (12.63 - 0.06 * 3) / 4
for i, (t, s) in enumerate(band):
    bx = 0.35 + i * (bw + 0.06)
    rect(s1, bx, 6.32, bw, 0.55)
    rect(s1, bx, 6.32, 0.05, 0.55, fill=PURPLE, line=None, rounded=False)
    multiline(s1, bx + 0.1, 6.32, bw - 0.12, 0.55, [
        (t, {"size": 8.4, "bold": True, "color": PURPLE_D, "sp": 0.95}),
        (s, {"size": 7, "color": MUTE, "sp": 0.95})], anchor=MSO_ANCHOR.MIDDLE)
footer(s1, 1)


# ============================================================ SLIDE 2
s2 = prs.slides.add_slide(BLANK)
header(s2, f"MODEL IN ACTION  ·  LIVE RUN, {cfg['risk_appetite'].upper()} APPETITE · {cfg['n_paths']:,} MONTE-CARLO PATHS",
       [("Across the basket the copilot cuts cost-risk ", {}),
        (f"~{opt['risk_cut_pct']:.0f}%", {"color": PURPLE}),
        (" at ~equal expected cost, and nets a further ", {}),
        (f"{port['diversification_benefit_pct']:.0f}%", {"color": PURPLE}),
        (" by managing commodities as one portfolio", {})])

# KPI strip
kpis = [("PORTFOLIO 6-MO SPEND", f"₹{port['total_spend_cr']:,.0f} Cr", "palm + crude + silver"),
        ("RISK IF SILOED", f"₹{port['additive_sd_cr']:,.0f} Cr", "risks added up"),
        ("RISK IF NETTED", f"₹{port['portfolio_sd_cr']:,.0f} Cr", f"-{port['diversification_benefit_pct']:.0f}% natural hedge"),
        ("FLAGSHIP RISK CUT", f"-{opt['risk_cut_pct']:.0f}%", "palm, vs do-nothing"),
        ("WORST-CASE TRIMMED", f"₹{dist['naive_car95']-dist['engine_car95']:,.0f} Cr", f"CaR95 {dist['naive_car95']:,.0f}→{dist['engine_car95']:,.0f}")]
kw = (12.63 - 0.08 * 4) / 5
for i, (l, v, n) in enumerate(kpis):
    kx = 0.35 + i * (kw + 0.08)
    rect(s2, kx, 1.2, kw, 0.78)
    multiline(s2, kx + 0.08, 1.2, kw - 0.14, 0.78, [
        (l, {"size": 6.8, "color": MUTE, "sp": 0.95}),
        (v, {"size": 15, "bold": True, "color": PURPLE_D, "sp": 0.95}),
        (n, {"size": 6.6, "color": MUTE, "sp": 0.95})], anchor=MSO_ANCHOR.MIDDLE)

# left column: per-commodity table + frontier
def acol(a): return AMBER if "PHASE" in a else (RED if "WAIT" in a else GREEN)
LX2, LW2 = 0.35, 6.1
rect(s2, LX2, 2.12, LW2, 1.62)
txt(s2, LX2 + 0.08, 2.16, LW2, 0.2, "PER-COMMODITY DECISIONS (THE WHOLE BASKET)", size=8.6, bold=True, color=PURPLE_D)
ct = s2.shapes.add_table(4, 6, Inches(LX2 + 0.08), Inches(2.38), Inches(LW2 - 0.16), Inches(0.86)).table
chd = ["Commodity", "Buy-now", "Lock/Opt/Float", "Risk cut", "Timing", "Act"]
for j, hh in enumerate(chd):
    ct.cell(0, j).text = hh
for i, (n, r) in enumerate(C.items(), 1):
    spl = r["optimise"]["split"]; a = r["score"]["action"]
    tv = "HEDGE" if r["timing"]["verdict"] == "HEDGE" else "PHYS"
    vals = [n.split("(")[0].strip(), a, f"{spl['locked_pct']:.0f}/{spl['option_pct']:.0f}/{spl['unhedged_pct']:.0f}",
            f"-{r['optimise']['risk_cut_pct']:.0f}%", tv, "AUTO" if r["act"]["decision"] == "AUTO-EXECUTE" else "ESC."]
    for j, v in enumerate(vals):
        ct.cell(i, j).text = v
for i in range(4):
    for j in range(6):
        c = ct.cell(i, j); c.fill.solid(); c.fill.fore_color.rgb = TINT if i == 0 else WHITE
        c.margin_left = Pt(3); c.margin_top = Pt(1); c.margin_bottom = Pt(1); c.vertical_anchor = MSO_ANCHOR.MIDDLE
        for pr in c.text_frame.paragraphs:
            for r in pr.runs:
                r.font.size = Pt(7.8); r.font.name = FONT; r.font.bold = (i == 0 or j == 0)
                r.font.color.rgb = PURPLE_D if i == 0 else INK
                if j == 1 and i > 0:
                    r.font.color.rgb = acol(list(C.values())[i-1]["score"]["action"]); r.font.bold = True
                if j == 4 and i > 0:
                    r.font.color.rgb = RED if list(C.values())[i-1]["timing"]["verdict"] == "HEDGE" else GREEN
                    r.font.bold = True
txt(s2, LX2 + 0.08, 3.40, LW2 - 0.16, 0.34,
    [("Timing", {"bold": True, "color": PURPLE_D, "size": 7.2}),
     (f" = DP schedules buys on the forward path incl. {P['timing']['holding_monthly_pct']:.2f}%/mo holding & {P['timing']['lead_weeks']}-wk lead, then best physical (buy-&-hold) vs hedge. ",
      {"size": 7.2, "color": MUTE}),
     ("Lock/Opt/Float", {"bold": True, "color": PURPLE_D, "size": 7.2}),
     (" adds to 100% (not \"60+50\").", {"size": 7.2, "color": MUTE})])

# frontier chart (native shapes)
rect(s2, LX2, 3.86, LW2, 2.36)
txt(s2, LX2 + 0.08, 3.9, LW2, 0.2, "DECIDE - EFFICIENT FRONTIER (PALM)", size=8.6, bold=True, color=PURPLE_D)
# plot box
px0, py0, pw, ph = LX2 + 0.55, 4.34, LW2 - 0.9, 1.5
Es = [g["E"] / 1e7 for g in opt["grid"]]; Ss = [g["sd"] / 1e7 for g in opt["grid"]]
xlo, xhi = min(Es) * 0.999, max(Es) * 1.001; yhi = max(Ss) * 1.1
def fx(v): return px0 + pw * (v - xlo) / (xhi - xlo)
def fy(v): return py0 + ph * (1 - v / yhi)
# axes
rect(s2, px0, py0, 0.012, ph, fill=INK, line=None, rounded=False)
rect(s2, px0, py0 + ph, pw, 0.012, fill=INK, line=None, rounded=False)
txt(s2, px0, py0 + ph + 0.03, pw, 0.2, "Expected cost (₹ Cr) → cheaper is left", size=6.8, color=MUTE, align=PP_ALIGN.CENTER)
txt(s2, px0 - 0.5, py0 - 0.18, 2.2, 0.16, "Risk = cost volatility (₹ Cr)", size=6.8, color=MUTE)
frontier = sorted(opt["frontier"], key=lambda g: g["E"])
# frontier connecting line via connectors
for k in range(len(frontier) - 1):
    ln = s2.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        Inches(fx(frontier[k]["E"]/1e7)), Inches(fy(frontier[k]["sd"]/1e7)),
        Inches(fx(frontier[k+1]["E"]/1e7)), Inches(fy(frontier[k+1]["sd"]/1e7)))
    ln.line.color.rgb = PURPLE; ln.line.width = Pt(1.5)
def dot(x, y, r, col):
    o = s2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - r), Inches(y - r), Inches(2*r), Inches(2*r))
    o.fill.solid(); o.fill.fore_color.rgb = col; o.line.fill.background(); o.shadow.inherit = False
for g in opt["grid"]:
    dot(fx(g["E"]/1e7), fy(g["sd"]/1e7), 0.035, MUTE)
for g in frontier:
    dot(fx(g["E"]/1e7), fy(g["sd"]/1e7), 0.045, PURPLE)
nv = next(g for g in opt["grid"] if g["cover"] == 0 and g["hedge"] == 0)
dot(fx(nv["E"]/1e7), fy(nv["sd"]/1e7), 0.06, RED)
txt(s2, fx(nv["E"]/1e7) + 0.08, fy(nv["sd"]/1e7) - 0.1, 1.6, 0.16, "do-nothing (all float)", size=6.5, color=RED)
dot(fx(ch["E"]/1e7), fy(ch["sd"]/1e7), 0.07, GREEN)
txt(s2, fx(ch["E"]/1e7) - 0.3, fy(ch["sd"]/1e7) + 0.06, 1.0, 0.16, "engine pick", size=6.8, bold=True, color=GREEN)

# right column
RX2, RW2 = 6.65, 6.33
# distribution chart (native area via freeform-ish bars)
rect(s2, RX2, 2.12, RW2, 2.15)
txt(s2, RX2 + 0.08, 2.16, RW2, 0.2, "THE PAY-OFF DISTRIBUTION (PALM, 20K FUTURES)", size=8.6, bold=True, color=PURPLE_D)
dx0, dy0, dw, dh = RX2 + 0.2, 2.5, RW2 - 0.4, 1.4
cx = dist["centers"]; nvd = dist["naive"]; end = dist["engine"]
ymax = max(max(nvd), max(end)) * 1.1
def ddx(i): return dx0 + dw * i / (len(cx) - 1)
def ddy(v): return dy0 + dh * (1 - v / ymax)
bw2 = dw / len(cx)
for i in range(len(cx)):
    if nvd[i] > 0:
        h = dh * nvd[i] / ymax
        b = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ddx(i)), Inches(dy0 + dh - h), Inches(bw2 * 0.9), Inches(h))
        b.fill.solid(); b.fill.fore_color.rgb = RED; b.fill.transparency = 0
        b.line.fill.background(); b.shadow.inherit = False
        b.fill.fore_color.rgb = RGBColor(0xF3, 0xC6, 0xCE)
for i in range(len(cx)):
    if end[i] > 0:
        h = dh * end[i] / ymax
        b = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ddx(i)), Inches(dy0 + dh - h), Inches(bw2 * 0.9), Inches(h))
        b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0xCF, 0xC2, 0xF3)
        b.line.fill.background(); b.shadow.inherit = False
rect(s2, dx0, dy0 + dh, dw, 0.01, fill=INK, line=None, rounded=False)
txt(s2, RX2 + 0.08, 3.9, RW2, 0.25, "Do-nothing (red) has a long expensive tail; engine plan (purple) trims the worst case at ~same expected cost.",
    size=7, color=MUTE)
# legend
rect(s2, RX2 + RW2 - 1.5, 2.34, 0.14, 0.12, fill=RGBColor(0xCF, 0xC2, 0xF3), line=None); txt(s2, RX2 + RW2 - 1.32, 2.3, 1.3, 0.16, "engine plan", size=6.6, color=INK)
rect(s2, RX2 + RW2 - 1.5, 2.5, 0.14, 0.12, fill=RGBColor(0xF3, 0xC6, 0xCE), line=None); txt(s2, RX2 + RW2 - 1.32, 2.46, 1.3, 0.16, "do-nothing", size=6.6, color=INK)

# insurance callout
rect(s2, RX2, 4.35, RW2, 0.72, fill=TINT, line=PURPLE)
multiline(s2, RX2 + 0.1, 4.35, RW2 - 0.2, 0.72, [
    (f"Read it as insurance, not a bet. Expected cost ~unchanged (₹{dist['naive_E']:,.0f}→₹{dist['engine_E']:,.0f} Cr) "
     f"but the 95%-worst-case cost falls ₹{dist['naive_car95']-dist['engine_car95']:,.0f} Cr. On real history the "
     f"cost line was {bt['vol_reduction_pct']:.0f}% steadier than naive buying.",
     {"size": 8.2, "color": INK, "sp": 1.03})], anchor=MSO_ANCHOR.MIDDLE)
# supply chain + governance
rect(s2, RX2, 5.15, RW2, 1.07)
txt(s2, RX2 + 0.08, 5.19, RW2, 0.2, "SUPPLY CHAIN & GOVERNANCE", size=8.6, bold=True, color=PURPLE_D)
multiline(s2, RX2 + 0.08, 5.4, RW2 - 0.16, 0.78, [
    (f"Supply first: nets inventory ({sup['horizon_demand']:,.0f}→{sup['net_procurement']:,.0f} {sup['unit']} to buy); "
     f"below reorder it forces a mandatory {sup['must_cover_now']:,.0f} {sup['unit']} continuity buy before price-optimising.",
     {"size": 7.8, "color": INK, "sp": 1.05}),
    (f"Council sets the preset (λ + caps); copilot auto-executes inside caps, escalates the rest "
     f"(palm ₹{act['immediate_order_value_cr']:.0f} Cr > ₹{act['auto_execute_cap_cr']} Cr cap → escalate). Full audit trail; override on.",
     {"size": 7.8, "color": MUTE, "sp": 1.05})], anchor=MSO_ANCHOR.TOP)

# bottom band
band2 = [("Cadence", "weekly re-solve + trigger alerts"),
         ("Decision rights", "auto in caps; council/CFO rest"),
         ("Proof", "backtest: steadier vs naive"),
         ("Scale-up", "1 commodity → basket → portfolio")]
for i, (t, s) in enumerate(band2):
    bx = 0.35 + i * (bw + 0.06)
    rect(s2, bx, 6.32, bw, 0.55)
    rect(s2, bx, 6.32, 0.05, 0.55, fill=PURPLE_D, line=None, rounded=False)
    multiline(s2, bx + 0.1, 6.32, bw - 0.12, 0.55, [
        (t, {"size": 8.2, "bold": True, "color": PURPLE_D, "sp": 0.95}),
        (s, {"size": 6.9, "color": MUTE, "sp": 0.95})], anchor=MSO_ANCHOR.MIDDLE)
footer(s2, 2)

out = os.path.join(HERE, "Commodity_Copilot_Slides_EDITABLE.pptx")
prs.save(out)
print("saved ->", out)
