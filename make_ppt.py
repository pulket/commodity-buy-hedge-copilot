"""
Turns copilot_slides.html into an actual PowerPoint (.pptx).
Each HTML slide is rendered to a crisp 2x image and placed full-bleed on a 16:9
slide - so the deck looks exactly like the HTML.
"""
import re, subprocess, os
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
TMP = "/private/tmp/claude-501/-Users-pulkitaggarwal-Desktop-kear/a68faf3d-ee66-4293-bd35-034622b178e8/scratchpad"
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

html = open(os.path.join(HERE, "copilot_slides.html")).read()
css = re.search(r"<style>(.*?)</style>", html, re.S).group(1)
slides = re.findall(r'<section class="slide">.*?</section>', html, re.S)
print(f"found {len(slides)} slides")

pngs = []
for i, sec in enumerate(slides, 1):
    one = (f"<!doctype html><html><head><meta charset='utf-8'><style>{css}\n"
           "html,body{margin:0;padding:0;background:#fff}.slide{box-shadow:none;margin:0}"
           f"</style></head><body>{sec}</body></html>")
    hp = os.path.join(TMP, f"slide{i}.html")
    pp = os.path.join(TMP, f"slide{i}.png")
    open(hp, "w").write(one)
    subprocess.run([CHROME, "--headless", "--disable-gpu", "--hide-scrollbars",
                    "--force-device-scale-factor=2", "--window-size=1280,720",
                    f"--screenshot={pp}", f"file://{hp}"],
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    pngs.append(pp)
    print(f"  slide {i} -> {pp}")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
for pp in pngs:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(pp, 0, 0, width=prs.slide_width, height=prs.slide_height)

out = os.path.join(HERE, "Commodity_Copilot_Slides.pptx")
prs.save(out)
print("saved ->", out)
